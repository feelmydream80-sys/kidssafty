from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re

ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
DARK = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
BLUE = RGBColor(0x34, 0x98, 0xDB)
GREEN = RGBColor(0x27, 0xAE, 0x60)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    if width is None: width = slide_width
    if height is None: height = Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_bullet_text(tf, text, size=13, color=DARK, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_after = Pt(4)
    return p

def parse_table_data(lines, idx):
    table_data = []
    while idx < len(lines):
        line = lines[idx].strip()
        if not line or line.startswith('#') or line.startswith('---'):
            break
        if line.startswith('|') and line.endswith('|'):
            cells = [c.strip() for c in line.split('|')[1:-1]]
            table_data.append(cells)
        idx += 1
    return table_data, idx

def add_table(slide, data, left, top, width, height):
    if len(data) < 2: return
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for ci in range(cols):
        header_cell = table.cell(0, ci)
        header_cell.text = data[0][ci]
        for p in header_cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
        header_cell.fill.solid()
        header_cell.fill.fore_color.rgb = BLUE

    for ri in range(1, rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            cell.text = data[ri][ci] if ci < len(data[ri]) else ''
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA) if ri % 2 == 0 else WHITE

    return table

def md_to_ppt(md_file, ppt_file):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw = prs.slide_width
    sh = prs.slide_height

    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    slides_raw = re.split(r'\n---\n', content)

    for slide_text in slides_raw:
        if not slide_text.strip():
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        lines = slide_text.strip().split('\n')

        title_match = re.search(r'# (.+)', slide_text)
        is_title_slide = title_match and ('표지' in title_match.group(1) or '맺음말' in title_match.group(1))

        if is_title_slide:
            add_bg(slide, ACCENT)
            title = title_match.group(1).replace('슬라이드 1: ', '').replace('슬라이드 15: ', '')
            add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.5), title, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

            sub_lines = [l.strip() for l in lines if l.strip() and not l.startswith('#') and not l.startswith('![') and not l.startswith('**')]
            y = Inches(3.8)
            for sl in sub_lines[:6]:
                add_textbox(slide, Inches(2), y, Inches(9), Inches(0.5), sl.replace('**', ''), size=16, color=WHITE, align=PP_ALIGN.CENTER)
                y += Inches(0.5)
            continue

        # 일반 슬라이드
        add_shape_bg(slide, ACCENT, top=0, height=Inches(1.1), width=sw)
        add_shape_bg(slide, RGBColor(0xC0, 0x39, 0x2B), top=Inches(1.1), height=Inches(0.06), width=sw)

        slide_num = ''
        title_text = ''
        if title_match:
            full = title_match.group(1)
            num_m = re.match(r'슬라이드\s*(\d+)', full)
            if num_m:
                slide_num = num_m.group(1)
                title_text = full.split(':', 1)[1].strip() if ':' in full else full
            else:
                title_text = full

            add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
                       f'{title_text}', size=24, bold=True, color=WHITE)

        # 본문 처리
        body_y = Inches(1.5)
        i = 0
        body_text = []
        table_start = False
        code_block = False

        while i < len(lines):
            line = lines[i]
            stripped = line.strip()

            if stripped.startswith('```'):
                code_block = not code_block
                i += 1
                continue

            if code_block:
                i += 1
                continue

            if stripped.startswith('#'):
                i += 1
                continue

            if stripped.startswith('|') and stripped.endswith('|'):
                table_rows = []
                while i < len(lines) and lines[i].strip().startswith('|'):
                    table_rows.append(lines[i].strip())
                    i += 1
                if len(table_rows) >= 3:
                    data = [c.strip() for c in table_rows[0].split('|')[1:-1]]
                    tdata = [data]
                    for tr in table_rows[2:]:
                        cells = [c.strip() for c in tr.split('|')[1:-1]]
                        if cells: tdata.append(cells)
                    tbl = add_table(slide, tdata, Inches(0.6), body_y, Inches(12), Inches(min(len(tdata) * 0.35, 4.5)))
                    body_y += Inches(len(tdata) * 0.35 + 0.3)
                continue

            # 체크박스 리스트
            if stripped.startswith('- [') or stripped.startswith('✅'):
                body_text.append(stripped)
                i += 1
                continue

            # 일반 텍스트
            if stripped and not stripped.startswith('!['):
                body_text.append(stripped)

            i += 1

        # 본문 텍스트 출력
        if body_text:
            tf = add_textbox(slide, Inches(0.6), body_y, Inches(12), Inches(5.5),
                            body_text[0] if body_text else '', size=14, color=DARK)
            for bt in body_text[1:]:
                if bt.startswith('  '):
                    add_bullet_text(tf, bt.strip(), size=13, color=GRAY, level=1)
                elif bt.startswith('- ['):
                    add_bullet_text(tf, bt, size=13, color=DARK)
                elif bt.startswith('✅') or bt.startswith('🎯'):
                    add_bullet_text(tf, bt, size=14, color=GREEN, bold=True)
                elif bt.startswith('**') and bt.endswith('**'):
                    add_bullet_text(tf, bt.replace('**', ''), size=14, color=ACCENT, bold=True)
                else:
                    add_bullet_text(tf, bt, size=13, color=DARK)

    prs.save(ppt_file)
    print(f'PPT 생성 완료: {ppt_file}')

if __name__ == '__main__':
    md_to_ppt('submission/ppt/KidsSafeAI_PPT_v4.md', 'submission/ppt/KidsSafeAI_PPT_v4.pptx')
